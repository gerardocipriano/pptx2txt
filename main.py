from pptx import Presentation
import glob
import os

out_dir = 'out'
if not os.path.exists(out_dir):
    os.makedirs(out_dir)

for eachfile in glob.glob("*.pptx"):
    prs = Presentation(eachfile)
    print(eachfile)
    print("----------------------")
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    if len(text_runs) >= 3 and all(not t for t in text_runs[-3:]):
                        text_runs = text_runs[:-2]
                    text_runs.append(text)
                else:
                    text_runs.append('')
    with open(os.path.join(out_dir, os.path.basename(eachfile) + '.txt'), 'w', encoding='utf-8') as f:
        f.write('\n'.join(text_runs))
